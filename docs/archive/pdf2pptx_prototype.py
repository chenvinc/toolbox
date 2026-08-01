# -*- coding: utf-8 -*-
"""
PDF -> PPTX 最终定版管线（一次到位，无需事后修补）
====================================================

目标形态（经多轮迭代确认的正确结构）：
  - 以 模版.pptx 为基底：继承其全部母版/主题/版式
  - 每一页使用「模版.pptx 第1页所用的版式」（即母版图片背景所在版式）
  - 页面上只有可编辑文字框，**没有任何图片、没有页面级背景覆盖**
  - 文字保留：真实字体名 / 字号(pt) / 颜色(RGB) / 粗体 / 斜体 / 精确坐标

历史教训（本脚本已全部规避）：
  1. [废弃] 整页PDF栅格图作背景 —— 栅格底图与实时渲染文字无法跨设备
     像素级重合，必然"错位"；且体积暴涨（43MB+）。本版完全不生成图片。
  2. [废弃] slide.background.fill.solid() 设白色兜底 —— 会在 slide XML 写入
     <p:bg> 页面级覆盖，把母版图片背景挡住（打开PPT只见白底）。
     本版绝不触碰 slide.background。
  3. [废弃] 使用"空白"版式 —— 与模版页不同母版/版式，观感不一致。
     本版直接取 模版.pptx 第1页的 slide_layout（自动定位，无需硬编码名字）。
  4. 坐标映射：模版 slide 尺寸(720x405.1pt) 恰与 PDF 页面 1:1，
     但本版仍按 (slide尺寸/PDF页面尺寸) 动态计算缩放，兼容其它尺寸的PDF。

用法:
  python pdf2pptx_final.py <pdf路径> [输出pptx路径] [模板路径]
  # 输出缺省 = 与pdf同名的.pptx；模板缺省 = 脚本同目录的 模版.pptx
批量:
  for f in *.pdf; do python pdf2pptx_final.py "$f"; done

依赖: pip install pymupdf python-pptx
"""
import os
import sys

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 字体归一
_FONT_ALIAS = {
    "MicrosoftYaHei": "Microsoft YaHei",
    "MicrosoftYaHeiUI": "Microsoft YaHei UI",
    "MicrosoftYaHei-Bold": "Microsoft YaHei",
    "MicrosoftYaHeiUI-Bold": "Microsoft YaHei UI",
    "MicrosoftYaHei-Light": "Microsoft YaHei Light",
    "ArialMT": "Arial",
    "Arial-BoldMT": "Arial",
    "Arial-ItalicMT": "Arial",
    "Arial-BoldItalicMT": "Arial",
}


def clean_font(raw):
    """去除子集前缀(ABCDEF+)与变体后缀(-Bold等)，返回可用字体家族名"""
    if not raw:
        return "Microsoft YaHei"
    name = raw.split("+", 1)[1] if "+" in raw else raw
    fam = name
    for suffix in ("-BoldItalic", "-Bold", "-Italic", "-Light", "-Regular"):
        if fam.endswith(suffix):
            fam = fam[: -len(suffix)]
            break
    return _FONT_ALIAS.get(fam, _FONT_ALIAS.get(name, fam)) or "Microsoft YaHei"


def is_bold(raw, flags):
    return "Bold" in (raw or "") or bool(flags & 16)


def is_italic(raw, flags):
    return "Italic" in (raw or "") or bool(flags & 2)


# ---------------------------------------------------------------- 模板处理
def get_reference_layout(prs):
    """取模板第1页所用的版式 —— 即用户认可的"这一页"的底子。
    必须在 strip_template_slides 之前调用。"""
    if len(prs.slides) == 0:
        raise RuntimeError("模板中没有幻灯片，无法确定参考版式")
    return prs.slides[0].slide_layout


def strip_template_slides(prs):
    """删除模板自带示例页，仅保留母版/主题/版式定义"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def remove_placeholders(slide):
    """删除 add_slide 从版式带进来的空占位符（模版原页上也没有这些形状）"""
    for ph in list(slide.placeholders):
        el = ph._element
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


# ---------------------------------------------------------------- 主转换
def convert(pdf_path, out_path, template_path):
    prs = Presentation(template_path)
    ref_layout = get_reference_layout(prs)     # 先取版式
    strip_template_slides(prs)                 # 再清示例页
    sw, sh = prs.slide_width, prs.slide_height

    doc = fitz.open(pdf_path)
    stats = {"pages": doc.page_count, "textboxes": 0, "runs": 0, "empty_pages": []}

    for pi in range(doc.page_count):
        page = doc[pi]
        # 动态缩放系数 (EMU per PDF point)；若尺寸恰好 1:1 则等价于 x12700
        ex = sw / page.rect.width
        ey = sh / page.rect.height

        slide = prs.slides.add_slide(ref_layout)
        remove_placeholders(slide)
        # !!! 关键：绝不设置 slide.background —— 保持无 <p:bg>，
        #     让版式/母版的图片背景自然透出。

        td = page.get_text("dict")
        page_has_text = False
        for b in td.get("blocks", []):
            if b.get("type") != 0:
                continue
            for ln in b.get("lines", []):
                spans = [s for s in ln.get("spans", []) if s.get("text")]
                if not spans:
                    continue
                xs0 = min(s["bbox"][0] for s in spans)
                ys0 = min(s["bbox"][1] for s in spans)
                xs1 = max(s["bbox"][2] for s in spans)
                ys1 = max(s["bbox"][3] for s in spans)

                tb = slide.shapes.add_textbox(
                    Emu(int(round(xs0 * ex))), Emu(int(round(ys0 * ey))),
                    Emu(max(1, int(round((xs1 - xs0) * ex)))),
                    Emu(max(1, int(round((ys1 - ys0) * ey)))),
                )
                tf = tb.text_frame
                tf.word_wrap = False
                tf.auto_size = None
                tf.vertical_anchor = MSO_ANCHOR.TOP
                tf.margin_left = tf.margin_right = Emu(0)
                tf.margin_top = tf.margin_bottom = Emu(0)

                p = tf.paragraphs[0]
                p.line_spacing = 1.0
                p.space_before = Pt(0)
                p.space_after = Pt(0)

                for sp in spans:
                    run = p.add_run()
                    run.text = sp["text"]
                    f = run.font
                    f.name = clean_font(sp.get("font"))
                    f.size = Pt(sp.get("size") or 12)
                    col = sp.get("color")
                    if col is not None:
                        c = col & 0xFFFFFF
                        f.color.rgb = RGBColor((c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF)
                    if is_bold(sp.get("font"), sp.get("flags", 0)):
                        f.bold = True
                    if is_italic(sp.get("font"), sp.get("flags", 0)):
                        f.italic = True
                    stats["runs"] += 1
                stats["textboxes"] += 1
                page_has_text = True

        if not page_has_text:
            stats["empty_pages"].append(pi + 1)   # 纯图片页(如封面)，源PDF本身无文字

        if (pi + 1) % 50 == 0:
            print("    ...%d/%d 页  文本框=%d" % (pi + 1, doc.page_count, stats["textboxes"]), flush=True)

    prs.save(out_path)
    doc.close()
    return stats


# ---------------------------------------------------------------- 自检
def self_check(out_path, pdf_path):
    """生成后自动校验三条铁律：无<p:bg>覆盖 / 无图片 / 页数一致"""
    import zipfile
    from lxml import etree
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    z = zipfile.ZipFile(out_path)
    slides = sorted(n for n in z.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    bad_bg, has_pic = [], []
    for sfile in slides:
        root = etree.fromstring(z.read(sfile))
        csld = root.find("{%s}cSld" % NS_P)
        if csld is not None and csld.find("{%s}bg" % NS_P) is not None:
            bad_bg.append(sfile)
        if b"<p:pic" in z.read(sfile):
            has_pic.append(sfile)
    n_pdf = fitz.open(pdf_path).page_count
    ok = (not bad_bg) and (not has_pic) and (len(slides) == n_pdf)
    print("    自检: 页数 %d/%d | 页面级bg覆盖 %d | 含图片页 %d  -> %s"
          % (len(slides), n_pdf, len(bad_bg), len(has_pic), "PASS" if ok else "FAIL"))
    return ok


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    pdf = os.path.abspath(sys.argv[1])
    out = os.path.abspath(sys.argv[2]) if len(sys.argv) > 2 else os.path.splitext(pdf)[0] + ".pptx"
    tpl = os.path.abspath(sys.argv[3]) if len(sys.argv) > 3 else \
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "模版.pptx")

    if not os.path.exists(tpl):
        print("错误: 找不到模板 %s" % tpl)
        sys.exit(2)

    print(">>> %s -> %s (模板: %s)" % (os.path.basename(pdf), os.path.basename(out), os.path.basename(tpl)), flush=True)
    st = convert(pdf, out, tpl)
    print("    统计: 页数=%d 文本框=%d 文字run=%d 无文字页=%s"
          % (st["pages"], st["textboxes"], st["runs"], st["empty_pages"] or "无"), flush=True)
    if not self_check(out, pdf):
        sys.exit(3)
