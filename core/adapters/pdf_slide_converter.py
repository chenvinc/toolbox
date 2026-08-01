"""PDF → PPTX 转换适配器（封装 pymupdf + python-pptx）。

实现 core/ports/io.PdfSlideConverter。逻辑逐行移植自定版管线
``docs/pdf2pptx_final.py``，规避其记载的全部历史教训：

1. 不生成任何整页栅格图 —— 栅格底图与实时渲染文字无法跨设备像素级重合，
   且体积暴涨；本适配器完全不产出图片。
2. 绝不触碰 ``slide.background`` —— 否则会在 slide XML 写入 ``<p:bg>``
   页面级覆盖，把母版图片背景挡住。
3. 不使用「空白」版式 —— 直接取模板第 1 页所用的 slide_layout
   （自动定位，无需硬编码版式名），保证观感与模板一致。
4. 坐标按 (slide 尺寸 / PDF 页面尺寸) 动态计算缩放，兼容任意尺寸 PDF。
"""
from __future__ import annotations

from typing import Any, Callable, Dict, List

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from core.ports.io import PdfSlideConverter
from shared.contracts import ConvertPdfRequest, ConvertPdfResult
from shared.errors import PdfReadError, TemplateInvalidError

# ---------------------------------------------------------------- 字体归一

_FONT_ALIAS: Dict[str, str] = {
    "MicrosoftYaHei": "Microsoft YaHei",
    "MicrosoftYaHeiUI": "Microsoft YaHei UI",
    "MicrosoftYaHei-Bold": "Microsoft YaHei",
    "MicrosoftYaHeiUI-Bold": "Microsoft YaHei UI",
    "MicrosoftYaHei-Light": "Microsoft YaHei Light",
    "ArialMT": "Arial",
    "Arial-BoldMT": "Arial",
    "Arial-ItalicMT": "Arial",
    "Arial-BoldItalicMT": "Arial",
}

_DEFAULT_FONT = "Microsoft YaHei"


def clean_font(raw: str | None) -> str:
    """去除子集前缀(ABCDEF+)与变体后缀(-Bold等)，返回可用字体家族名。"""
    if not raw:
        return _DEFAULT_FONT
    name = raw.split("+", 1)[1] if "+" in raw else raw
    fam = name
    for suffix in ("-BoldItalic", "-Bold", "-Italic", "-Light", "-Regular"):
        if fam.endswith(suffix):
            fam = fam[: -len(suffix)]
            break
    return _FONT_ALIAS.get(fam, _FONT_ALIAS.get(name, fam)) or _DEFAULT_FONT


def is_bold(raw: str | None, flags: int) -> bool:
    """按字体名与 span flags（bit4）判断粗体。"""
    return "Bold" in (raw or "") or bool(flags & 16)


def is_italic(raw: str | None, flags: int) -> bool:
    """按字体名与 span flags（bit1）判断斜体。"""
    return "Italic" in (raw or "") or bool(flags & 2)


# ---------------------------------------------------------------- 模板处理

def _get_reference_layout(prs: Any) -> Any:
    """取模板第 1 页所用版式（即用户认可页面的底子）。

    必须在 ``_strip_template_slides`` 之前调用。
    """
    if len(prs.slides) == 0:
        raise TemplateInvalidError("模板中没有幻灯片，无法确定参考版式")
    return prs.slides[0].slide_layout


def _strip_template_slides(prs: Any) -> None:
    """删除模板自带示例页，仅保留母版/主题/版式定义。"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:  # noqa: BLE001 - rel 缺失时跳过即可
            pass


def _remove_placeholders(slide: Any) -> None:
    """删除 add_slide 从版式带进来的空占位符（模板原页上没有这些形状）。"""
    for ph in list(slide.placeholders):
        el = ph._element
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


# ---------------------------------------------------------------- 适配器

class PdfSlideConverterAdapter(PdfSlideConverter):
    """按定版管线把 PDF 文字逐页搬进模板版式的可编辑文本框。"""

    def convert(
        self,
        request: ConvertPdfRequest,
        on_progress: Callable[[int, int], None],
    ) -> ConvertPdfResult:
        try:
            prs = Presentation(request.template_path)
        except Exception as exc:
            raise TemplateInvalidError(f"模板打开失败：{exc}") from exc

        ref_layout = _get_reference_layout(prs)   # 先取版式
        _strip_template_slides(prs)               # 再清示例页
        sw, sh = prs.slide_width, prs.slide_height
        if not sw or not sh:
            raise TemplateInvalidError("模板幻灯片尺寸缺失，无法计算坐标映射")

        try:
            doc = fitz.open(request.pdf_path)
        except Exception as exc:
            raise PdfReadError(f"PDF 打开失败：{exc}") from exc

        textboxes = 0
        runs = 0
        empty_pages: List[int] = []

        try:
            total = int(doc.page_count)
            for pi in range(total):
                page = doc[pi]
                # 动态缩放系数 (EMU per PDF point)；尺寸恰 1:1 时等价于 x12700
                ex = float(sw) / float(page.rect.width)
                ey = float(sh) / float(page.rect.height)

                slide = prs.slides.add_slide(ref_layout)
                _remove_placeholders(slide)
                # !!! 关键：绝不设置 slide.background —— 保持无 <p:bg>，
                #     让版式/母版的图片背景自然透出。

                td = page.get_text("dict")
                page_has_text = False
                for block in td.get("blocks", []):
                    if block.get("type") != 0:
                        continue
                    for line in block.get("lines", []):
                        spans = [s for s in line.get("spans", []) if s.get("text")]
                        if not spans:
                            continue
                        added_runs = self._add_line_textbox(slide, spans, ex, ey)
                        runs += added_runs
                        textboxes += 1
                        page_has_text = True

                if not page_has_text:
                    empty_pages.append(pi + 1)  # 纯图片页(如封面)，源 PDF 无文字

                on_progress(pi + 1, total)

            prs.save(request.output_path)
        finally:
            doc.close()

        return ConvertPdfResult(
            output_path=request.output_path,
            page_count=total,
            textbox_count=textboxes,
            run_count=runs,
            empty_pages=empty_pages,
        )

    @staticmethod
    def _add_line_textbox(
        slide: Any, spans: List[Dict[str, Any]], ex: float, ey: float
    ) -> int:
        """为一行 spans 添加一个精确坐标文本框，返回写入的 run 数。"""
        xs0 = min(float(s["bbox"][0]) for s in spans)
        ys0 = min(float(s["bbox"][1]) for s in spans)
        xs1 = max(float(s["bbox"][2]) for s in spans)
        ys1 = max(float(s["bbox"][3]) for s in spans)

        tb = slide.shapes.add_textbox(
            Emu(int(round(xs0 * ex))), Emu(int(round(ys0 * ey))),
            Emu(max(1, int(round((xs1 - xs0) * ex)))),
            Emu(max(1, int(round((ys1 - ys0) * ey)))),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        count = 0
        for sp in spans:
            run = p.add_run()
            run.text = sp["text"]
            f = run.font
            f.name = clean_font(sp.get("font"))
            f.size = Pt(sp.get("size") or 12)
            col = sp.get("color")
            if col is not None:
                c = int(col) & 0xFFFFFF
                f.color.rgb = RGBColor(  # type: ignore[no-untyped-call]
                    (c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF
                )
            if is_bold(sp.get("font"), sp.get("flags", 0)):
                f.bold = True
            if is_italic(sp.get("font"), sp.get("flags", 0)):
                f.italic = True
            count += 1
        return count
