"""PPT 生成适配器（封装 python-pptx）。

实现 core/ports/io.PptxWriter。把原 utils.generate_pptx 的纯生成逻辑迁移至此，
由 PptxServiceImpl 通过端口调用，外部（开发/测试）可用 FakePptxWriter 替换。
"""
from __future__ import annotations

from typing import Callable, List, TYPE_CHECKING

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.exc import InvalidXmlError

from core.ports.io import PptxWriter

if TYPE_CHECKING:
    from docx.text.run import Run
    from pptx.presentation import Presentation as PptxPresentation
    from pptx.slide import SlideLayout


def _set_font(run: "Run", font_name: str, font_size: int) -> None:
    """为文本 run 设置字体名称、字号以及东亚字体。"""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {})
    ea.set("typeface", font_name)
    rPr.append(ea)


def _remove_first_slide(prs: "PptxPresentation") -> "SlideLayout":
    """删除演示文稿的第一张幻灯片，返回其 slide_layout 供后续复用。

    加固：rId 缺失或 python-pptx 移除私有属性时优雅降级，保证生成不崩溃。
    """
    if len(prs.slides) == 0:
        return prs.slide_layouts[-1]
    try:
        layout: "SlideLayout" = prs.slides[0].slide_layout
    except InvalidXmlError:
        return prs.slide_layouts[-1]
    sldId_lst = getattr(prs.slides, "_sldIdLst", None)
    if sldId_lst is None:
        return layout
    try:
        sldId = sldId_lst[0]
        rId = sldId.get(qn("r:id"))
        if rId is not None:
            prs.part.drop_rel(rId)
        sldId_lst.remove(sldId)
    except (KeyError, TypeError, InvalidXmlError):
        pass
    return layout


class PptxWriterAdapter(PptxWriter):
    """基于 PPT 模板为每道题生成两页幻灯片并保存。"""

    def build(
        self,
        template_path: str,
        questions: List[List[str]],
        font_name: str,
        font_size: int,
        output_path: str,
        line_spacing: float,
        first_line_indent: bool,
        on_progress: Callable[[int, int], None],
    ) -> int:
        prs = Presentation(template_path)
        slide_layout = _remove_first_slide(prs)
        total = len(questions)
        for qi, question in enumerate(questions):
            if on_progress is not None:
                on_progress(qi + 1, total)
            for _ in range(2):
                slide = prs.slides.add_slide(slide_layout)
                for shape in list(slide.shapes):
                    sp = shape._element
                    sp.getparent().remove(sp)

                sw = prs.slide_width or 0
                sh = prs.slide_height or 0
                textbox_w = int(int(sw) * 0.8)
                textbox_l = int(int(sw) * 0.1)
                textbox_h = int(int(sh) * 0.75)
                textbox_t = int(int(sh) * 0.15)

                textbox = slide.shapes.add_textbox(textbox_l, textbox_t, textbox_w, textbox_h)
                tf = textbox.text_frame
                tf.word_wrap = True

                for i, para_text in enumerate(question):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    if first_line_indent:
                        para_text = "　　" + para_text
                    run = p.add_run()
                    run.text = para_text
                    _set_font(run, font_name, font_size)
                    p.space_after = Pt(0)
                    p.space_before = Pt(0)
                    p.line_spacing = line_spacing

        prs.save(output_path)
        return total * 2
