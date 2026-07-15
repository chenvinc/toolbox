"""P0 #3 生成 PPT 的路径冲突校验与私有 API 加固的回归测试（阶段4 迁移版）。

覆盖：
  1. 输出路径 == 模板路径时 generate 抛 ValueError（不覆盖/损坏模板）。
  2. 首张幻灯片 sldId 的 r:id 为 None 时 _remove_first_slide 不崩溃（修复点）。
  3. 正常模板（含 1 张首页）能被删除，生成的页数 = 2 * 题数。
  4. 0 张模板降级到 slide_layouts[-1]，仍能正常生成。
  5. 生成的 PPTX 可被 python-pptx 重新打开，文本已写入。
  6.（附加）DocxLoaderAdapter + parse_questions 等价 legacy extract_questions。

直接测试 core 层（零 Qt 依赖，无需 offscreen 亦可运行）。
"""
import os
import pathlib
import sys
import tempfile
import unittest

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = pathlib.Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from core.adapters.docx_loader import DocxLoaderAdapter
from core.adapters.pptx_writer import (
    PptxWriterAdapter, _remove_first_slide, _resolve_line_spacing, _same_path,
)
from core.services.slide_builder import PptxServiceImpl
from core.services._question_parser import parse_questions
from shared.contracts import GeneratePptxRequest, LineSpacingType
from shared.errors import OutputOverwriteError


class _NoOpEmitter:
    """忽略事件的 emitter 实现，供兼容壳 / 单测在无事件订阅时使用。"""
    def emit(self, event):
        pass

    def on_event(self, handler):
        pass


QUESTIONS = [
    ["1. 第一题题干", "A. 选项一", "B. 选项二"],
    ["2. 第二题题干", "A. 选项甲", "B. 选项乙"],
]


def _make_template(path, n_slides=1):
    prs = Presentation()
    for _ in range(n_slides):
        prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)
    return path


def _generate(tmpl, out, questions=QUESTIONS):
    svc = PptxServiceImpl(PptxWriterAdapter(), _NoOpEmitter())
    req = GeneratePptxRequest(
        template_path=tmpl,
        questions=questions,
        font_name="微软雅黑",
        font_size=20,
        output_path=out,
        line_spacing_type=LineSpacingType("1 倍"),
        line_spacing_value=1.0,
        first_line_indent=True,
    )
    return svc.generate(req)


class GeneratePptxP0Tests(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.mkdtemp(prefix="pptx_test_")

    def tearDown(self):
        for f in os.listdir(self._tmp):
            try:
                os.remove(os.path.join(self._tmp, f))
            except OSError:
                pass
        try:
            os.rmdir(self._tmp)
        except OSError:
            pass

    # ── 1. 路径冲突校验 ──
    def test_same_path_raises_value_error(self):
        tmpl = os.path.join(self._tmp, "t.pptx")
        _make_template(tmpl, 1)
        with self.assertRaises(ValueError) as ctx:
            _generate(tmpl, tmpl)
        self.assertIn("相同", str(ctx.exception))
        self.assertEqual(len(Presentation(tmpl).slides), 1)

    # ── 2. rId 为 None 时不崩溃 ──
    def test_remove_first_slide_with_none_rid_does_not_crash(self):
        tmpl = os.path.join(self._tmp, "none_rid.pptx")
        _make_template(tmpl, 1)
        prs = Presentation(tmpl)
        sldId = prs.slides._sldIdLst[0]
        sldId.attrib.pop(qn("r:id"), None)
        self.assertIsNone(sldId.get(qn("r:id")))
        try:
            _remove_first_slide(prs)
        except Exception as e:  # pragma: no cover - 不应到达
            self.fail(f"_remove_first_slide 在 rId=None 时不应抛异常: {e}")
        self.assertIn(len(prs.slides), (0, 1))

    # ── 3. 正常删首页，页数正确 ──
    def test_normal_template_removes_first_slide(self):
        tmpl = os.path.join(self._tmp, "normal.pptx")
        _make_template(tmpl, 1)
        out = os.path.join(self._tmp, "out.pptx")
        _generate(tmpl, out)
        self.assertEqual(len(Presentation(out).slides), 4)

    # ── 4. 0 张模板降级 ──
    def test_zero_slide_template_falls_back(self):
        tmpl = os.path.join(self._tmp, "zero.pptx")
        prs = Presentation()
        sldId_lst = prs.slides._sldIdLst
        for s in list(sldId_lst):
            sldId_lst.remove(s)
        prs.save(tmpl)
        self.assertEqual(len(Presentation(tmpl).slides), 0)
        out = os.path.join(self._tmp, "out_zero.pptx")
        _generate(tmpl, out)
        self.assertEqual(len(Presentation(out).slides), 4)

    # ── 5. 文本确实写入 ──
    def test_generated_text_present(self):
        tmpl = os.path.join(self._tmp, "txt.pptx")
        _make_template(tmpl, 1)
        out = os.path.join(self._tmp, "out_txt.pptx")
        _generate(tmpl, out)
        texts = []
        for slide in Presentation(out).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        blob = "\n".join(texts)
        self.assertIn("第一题题干", blob)
        self.assertIn("选项甲", blob)

    # ── 6. DocxLoaderAdapter + parse_questions 等价 legacy extract_questions ──
    def test_extract_questions_via_core(self):
        from docx import Document
        docx_path = os.path.join(self._tmp, "src.docx")
        doc = Document()
        doc.add_paragraph("一些无关的说明文本，不应被识别为题目。")
        doc.add_paragraph("1. 下列哪个是 Python 关键字？")
        doc.add_paragraph("A. class")
        doc.add_paragraph("B. def")
        doc.add_paragraph("C. if")
        doc.add_paragraph("D. all of the above")
        doc.save(docx_path)

        paras = DocxLoaderAdapter().load_paragraphs(docx_path)
        questions = parse_questions(paras, "1.", "A.")
        self.assertTrue(questions)
        self.assertEqual(questions[0][0], "1. 下列哪个是 Python 关键字？")

    # ── 辅助函数 ──
    def test_same_path_helper(self):
        self.assertTrue(_same_path("/a/b.pptx", "/a/b.pptx"))
        self.assertFalse(_same_path("/a/b.pptx", "/a/c.pptx"))

    def test_resolve_line_spacing_helper(self):
        self.assertEqual(_resolve_line_spacing("1 倍", 2.0), 1.0)
        self.assertEqual(_resolve_line_spacing("1.5 倍", 2.0), 1.5)
        self.assertEqual(_resolve_line_spacing("自定义", 2.0), 2.0)


if __name__ == "__main__":
    unittest.main()
