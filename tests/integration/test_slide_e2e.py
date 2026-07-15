"""Word → Slide 端到端测试（真实适配器，无 Qt）。

用真实的 DocxLoaderAdapter / PptxWriterAdapter 走完整 core 栈：
  真实 .docx → extract → 真实 .pptx 生成，验证契约、服务与适配器协同正确。
"""
import os
import tempfile
import unittest

from docx import Document
from pptx import Presentation

from shared.contracts import ExtractQuestionsRequest, GeneratePptxRequest
from core.di import Container


class CollectingEmitter:
    def __init__(self):
        self.events = []
        self._handlers = []

    def emit(self, event):
        self.events.append(event)
        for h in self._handlers:
            h(event)

    def on_event(self, handler):
        self._handlers.append(handler)


class SyncTaskRunner:
    def submit(self, func, *, args=(), kwargs=None, on_progress=None,
               on_result=None, on_error=None):
        try:
            result = func(*args, **(kwargs or {}))
            if on_result:
                on_result(result)
            return None
        except Exception as exc:
            if on_error:
                on_error(exc)
            return None


class SlideE2ETests(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.mkdtemp(prefix="slide_e2e_")

    def _make_docx(self, path):
        doc = Document()
        for line in ["1. 题目一", "A. 选项一", "B. 选项二",
                     "2. 题目二", "A. a", "B. b"]:
            doc.add_paragraph(line)
        doc.save(path)

    def _make_template(self, path, n=1):
        prs = Presentation()
        for _ in range(n):
            prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(path)

    def test_extract_and_generate_real_adapters(self):
        docx = os.path.join(self._tmp, "in.docx")
        self._make_docx(docx)
        tmpl = os.path.join(self._tmp, "t.pptx")
        self._make_template(tmpl, 1)
        out = os.path.join(self._tmp, "out.pptx")

        c = Container.build(task_runner=SyncTaskRunner(), event_emitter=CollectingEmitter())

        res = c.resolve("extraction").extract(ExtractQuestionsRequest(doc_path=docx))
        self.assertEqual(len(res.questions), 2)

        gres = c.resolve("pptx").generate(
            GeneratePptxRequest(
                template_path=tmpl,
                questions=res.questions,
                font_name="Arial", font_size=18, output_path=out,
            )
        )
        self.assertTrue(os.path.exists(out))
        self.assertEqual(gres.page_count, 4)
        self.assertEqual(len(Presentation(out).slides), 4)


if __name__ == "__main__":
    unittest.main()
