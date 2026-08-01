"""PdfSlideConverterAdapter 端到端测试（真实 pymupdf / python-pptx，临时文件）。

验证定版管线三条铁律（对应 docs/pdf2pptx_final.py 的自检）：
  1. 输出无页面级 <p:bg> 覆盖（母版背景不被遮挡）；
  2. 输出不含任何 <p:pic> 图片；
  3. 输出页数与源 PDF 一致。
另验证文字内容/无文字页统计与字体归一纯函数。
"""
import os
import tempfile
import unittest
import zipfile

import fitz
from pptx import Presentation
from pptx.util import Inches

from core.adapters.pdf_slide_converter import (
    PdfSlideConverterAdapter, clean_font, is_bold, is_italic,
)
from shared.contracts import ConvertPdfRequest
from shared.errors import PdfReadError, TemplateInvalidError


def _make_pdf(path, texts_per_page):
    """生成简单多页 PDF；texts_per_page 中空列表表示无文字页。"""
    doc = fitz.open()
    for texts in texts_per_page:
        page = doc.new_page(width=720, height=405)
        for i, text in enumerate(texts):
            page.insert_text((72, 72 + i * 30), text, fontsize=18)
    doc.save(path)
    doc.close()


def _make_template(path):
    """生成含 1 页示例幻灯片的最小 PPT 模板。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    prs.save(path)


class PdfSlideConverterE2ETests(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self._tmp.cleanup)
        self.pdf_path = os.path.join(self._tmp.name, "in.pdf")
        self.tpl_path = os.path.join(self._tmp.name, "tpl.pptx")
        self.out_path = os.path.join(self._tmp.name, "out.pptx")
        _make_pdf(self.pdf_path, [["Hello PDF", "Second line"], [], ["Page three"]])
        _make_template(self.tpl_path)
        self.adapter = PdfSlideConverterAdapter()

    def _convert(self):
        progress = []
        result = self.adapter.convert(
            ConvertPdfRequest(
                pdf_path=self.pdf_path,
                template_path=self.tpl_path,
                output_path=self.out_path,
            ),
            lambda c, t: progress.append((c, t)),
        )
        return result, progress

    def test_page_count_matches_pdf(self):
        result, progress = self._convert()
        self.assertEqual(result.page_count, 3)
        self.assertEqual(len(Presentation(self.out_path).slides), 3)
        self.assertEqual(progress, [(1, 3), (2, 3), (3, 3)])

    def test_empty_pages_reported(self):
        result, _ = self._convert()
        self.assertEqual(result.empty_pages, [2])
        self.assertGreater(result.textbox_count, 0)
        self.assertGreater(result.run_count, 0)

    def test_no_page_background_and_no_pictures(self):
        self._convert()
        with zipfile.ZipFile(self.out_path) as z:
            slides = [
                n for n in z.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            ]
            self.assertEqual(len(slides), 3)
            for name in slides:
                xml = z.read(name)
                self.assertNotIn(b"<p:bg", xml, f"{name} 不应有页面级背景覆盖")
                self.assertNotIn(b"<p:pic", xml, f"{name} 不应含图片")

    def test_text_content_preserved(self):
        self._convert()
        prs = Presentation(self.out_path)
        texts = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        joined = "\n".join(texts)
        self.assertIn("Hello PDF", joined)
        self.assertIn("Page three", joined)

    def test_invalid_pdf_raises(self):
        bad = os.path.join(self._tmp.name, "bad.pdf")
        with open(bad, "wb") as f:
            f.write(b"not a pdf")
        with self.assertRaises(PdfReadError):
            self.adapter.convert(
                ConvertPdfRequest(
                    pdf_path=bad, template_path=self.tpl_path,
                    output_path=self.out_path,
                ),
                lambda c, t: None,
            )

    def test_template_without_slides_raises(self):
        empty_tpl = os.path.join(self._tmp.name, "empty.pptx")
        Presentation().save(empty_tpl)  # 无任何幻灯片
        with self.assertRaises(TemplateInvalidError):
            self.adapter.convert(
                ConvertPdfRequest(
                    pdf_path=self.pdf_path, template_path=empty_tpl,
                    output_path=self.out_path,
                ),
                lambda c, t: None,
            )


class FontHelperTests(unittest.TestCase):
    def test_clean_font_strips_subset_prefix_and_suffix(self):
        self.assertEqual(clean_font("ABCDEF+MicrosoftYaHei-Bold"), "Microsoft YaHei")
        self.assertEqual(clean_font("ArialMT"), "Arial")
        self.assertEqual(clean_font("SomeFont-Italic"), "SomeFont")

    def test_clean_font_empty_falls_back(self):
        self.assertEqual(clean_font(None), "Microsoft YaHei")
        self.assertEqual(clean_font(""), "Microsoft YaHei")

    def test_bold_italic_flags(self):
        self.assertTrue(is_bold("Foo-Bold", 0))
        self.assertTrue(is_bold("Foo", 16))
        self.assertFalse(is_bold("Foo", 0))
        self.assertTrue(is_italic("Foo-Italic", 0))
        self.assertTrue(is_italic("Foo", 2))
        self.assertFalse(is_italic("Foo", 0))


if __name__ == "__main__":
    unittest.main()
