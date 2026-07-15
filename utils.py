"""题目文本提取工具 — 从 Word 文档中解析题目和选项。"""

from __future__ import annotations

import os
import re
from typing import Callable, List, Optional, Tuple
from typing import TYPE_CHECKING

from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.exc import InvalidXmlError

if TYPE_CHECKING:
    from docx.text.run import Run
    from pptx.presentation import Presentation as PptxPresentation
    from pptx.slide import SlideLayout


# ---------------------------------------------------------------------------
# 模块级常量
# ---------------------------------------------------------------------------

# 题号数值合理区间（内置宽泛匹配时生效）
_VALID_QUESTION_NUM_MIN = 1
_VALID_QUESTION_NUM_MAX = 200

# 常见中文题号分隔符（.．。、，, 全角/半角点号与顿号）
_COMMON_NUM_SEPS = r'.．。、，,'


# ---------------------------------------------------------------------------
# 辅助：题号 / 选项正则构建
# ---------------------------------------------------------------------------

def _is_raw_regex(pattern: str) -> bool:
    """判断字符串是否已经是正则表达式（而非格式示例）。"""
    return bool(re.search(r'\\[dDwWsS]|[+*?\[\]{}^|()]', pattern))


def _build_num_regex(num_pattern: str) -> Tuple[Optional[str], bool]:
    """根据用户输入的题号格式构建正则表达式，返回 (regex_str, needs_validation)。

    三种模式：
      1. 空/None   → 返回 None, True（调用方切换内置宽泛正则）
      2. 原始正则  → 返回 pattern, False（检测到 \\d、[]、+ 等特殊字符）
      3. 格式示例  → 返回 \"\\d+\\s*[...]+\"，兼容常见分隔符变体

    格式示例说明：
      用户输入 "1." → 题号形如 "36." / "36．" / "36、" 均能匹配
      用户输入 "1、" → 同上（分隔符字符类统一覆盖点号与顿号家族）
    """
    if not num_pattern or not num_pattern.strip():
        return None, True

    if _is_raw_regex(num_pattern):
        return num_pattern, False

    # 格式示例：提取用户输入中的分隔符部分（去除前导数字）
    m = re.match(r'\d+', num_pattern)
    if m:
        sep = num_pattern[m.end():]
        if sep and all(c in _COMMON_NUM_SEPS for c in sep):
            sep_class = _COMMON_NUM_SEPS + r'\-/～~·'
        else:
            sep_class = re.escape(sep) if sep else _COMMON_NUM_SEPS
        return r'\d+\s*[' + sep_class + r']+', False

    return re.escape(num_pattern), False


def _build_opt_regex(opt_prefix: str) -> Optional[str]:
    """根据用户输入的选项前缀构建正则表达式。

    - 空/None → 返回内置宽泛选项正则（兼容多种格式）
    - "A." 形式 → 扩展为 [A-Z]\\. 以匹配任意选项字母
    - 原始正则 → 原样返回
    - 其他 → 作为字面量使用
    """
    if not opt_prefix or not opt_prefix.strip():
        return None

    if _is_raw_regex(opt_prefix):
        return opt_prefix

    if re.fullmatch(r'[A-Za-z]\W', opt_prefix):
        return '[A-Za-z]' + re.escape(opt_prefix[1:])

    return re.escape(opt_prefix)


# ---------------------------------------------------------------------------
# 内置宽泛正则（兜底：用户未指定格式时启用）
# ---------------------------------------------------------------------------

# 内置宽泛题号正则
# group(1) 捕获题号原始字符串（可能含数字间空格），供数值校验
_BROAD_NUM_RE = re.compile(
    r'^\s*'                          # 行首空白
    r'(\d(?:\s*\d)*)'                # 题号数字（允许空格穿插）
    r'\s*'                           # 数字与分隔符之间的空白
    r'[.。,，、．\-/／～~—·⦁・･]+'      # 一个或多个分隔符
    r'\s*'                           # 分隔符后的空白
)

# 内置宽泛选项正则（兜底）
# 兼容格式：A.  A、  A．  (A)  （A）  a.  A ) 等，限定 A-D（大小写均可）
_BROAD_OPT_RE = re.compile(
    r'(?:[(（]\s*)?'                 # 可选前括号
    r'[A-Da-d]'                      # 选项字母 A-D
    r'(?:\s*[.。,，、．)）])'           # 必需：分隔符
    r'(?:\s*[)）])?'                 # 可选后括号
)


def _is_valid_question_num(num_str: str) -> bool:
    """校验题号是否在合理范围内（宽泛匹配模式专用）。"""
    try:
        n = int(re.sub(r'\s+', '', num_str))
    except ValueError:
        return False
    return _VALID_QUESTION_NUM_MIN <= n <= _VALID_QUESTION_NUM_MAX


# ---------------------------------------------------------------------------
# 行内选项拆分 & 题目收尾
# ---------------------------------------------------------------------------

def _split_inline_options(lines: List[str], opt_re: "re.Pattern[str]") -> List[str]:
    """将同一行内包含多个选项的文本拆分为多行。

    遍历每一行，如果一行中包含多个选项前缀匹配（如 "A. xxx B. xxx"），
    则将其按匹配位置拆分为独立的行；如果某行在第一个选项前还有题干文本，
    则题干文本也会被单独提取为一行。
    """
    result = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            result.append(line)
            continue
        matches = list(opt_re.finditer(stripped))
        if len(matches) <= 1:
            result.append(line)
        else:
            first_start = matches[0].start()
            if first_start > 0:
                prefix = stripped[:first_start].strip()
                if prefix:
                    result.append(prefix)
            for i, m in enumerate(matches):
                start = m.start()
                if i + 1 < len(matches):
                    end = matches[i + 1].start()
                    result.append(stripped[start:end].rstrip())
                else:
                    result.append(stripped[start:].rstrip())
    return result


def _finish_question(lines: List[str], opt_re: "re.Pattern[str]") -> Optional[List[str]]:
    """完成一道题目的收集，并格式化为行列表。

    双重校验：
      1. 题干（首行）非空
      2. 至少 2 个有效选项
    通过后调用 _split_inline_options 处理行内多选项。
    """
    if not lines or not lines[0].strip():
        return None
    full_text = '\n'.join(lines)
    if len(opt_re.findall(full_text)) < 2:
        return None
    return _split_inline_options(lines, opt_re)


# ---------------------------------------------------------------------------
# 公开接口
# ---------------------------------------------------------------------------

def extract_questions(doc_path: str, num_pattern: str, opt_prefix: str) -> List[List[str]]:
    """从 Word 文档中提取所有题目。

    Args:
        doc_path:    文档路径
        num_pattern: 题号格式示例（如 "1."）或原始正则；空字符串时启用内置宽泛匹配
        opt_prefix:  选项前缀示例（如 "A."）或原始正则；空字符串时启用内置宽泛匹配

    Returns:
        List[List[str]] — 每道题为一个行列表（题干 + 选项行）

    格式示例模式：
      用户传入 "1." 即可同时匹配 "36."、"36．"、"36、" 等常见中文试卷格式，
      无需每次根据文档格式手动调整题号分隔符。

    原始正则模式：
      若传入 \\d、[]、+ 等正则特殊字符，则直接作为正则表达式使用。

    内置宽泛匹配：
      传入空字符串时，使用内置正则兜底，可兼容更多异常格式（数字空格、
      重复分隔符、多种括号选项等），并校验题号在 1-200 范围内。
    """
    # ── 构建题号正则 ──
    num_re_str, needs_validation = _build_num_regex(num_pattern)
    if num_re_str is None:
        num_re = _BROAD_NUM_RE
        needs_validation = True
    else:
        num_re = re.compile(r'^\s*' + num_re_str)

    # ── 构建选项正则 ──
    opt_re_str = _build_opt_regex(opt_prefix)
    if opt_re_str is None:
        opt_re = _BROAD_OPT_RE
    else:
        opt_re = re.compile(opt_re_str)

    # ── 逐段扫描 ──
    doc = Document(doc_path)
    questions: List[List[str]] = []
    current_lines: List[str] = []
    collecting = False

    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            continue

        m = num_re.match(text)
        is_num_start = m is not None and (
            not needs_validation or _is_valid_question_num(m.group(1))
        )

        if is_num_start:
            if collecting and current_lines:
                q = _finish_question(current_lines, opt_re)
                if q is not None:
                    questions.append(q)
            current_lines = [text]
            collecting = True
        elif collecting:
            current_lines.append(text)

    if collecting and current_lines:
        q = _finish_question(current_lines, opt_re)
        if q is not None:
            questions.append(q)

    return questions


# ---------------------------------------------------------------------------
# PPT 生成工具（无需修改）
# ---------------------------------------------------------------------------

def _set_font(run: "Run", font_name: str, font_size: int) -> None:
    """为文本 run 设置字体名称、字号以及东亚字体。"""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {})
    ea.set('typeface', font_name)
    rPr.append(ea)


def _resolve_line_spacing(line_spacing_type: str, line_spacing_value: float) -> float:
    """根据行间距类型解析为实际的行间距数值。"""
    if line_spacing_type == "1.5 倍":
        return 1.5
    if line_spacing_type == "自定义":
        return line_spacing_value
    return 1.0


def _same_path(a: str, b: str) -> bool:
    """判断两个路径是否指向同一文件（规范化大小写与绝对路径后比较）。"""
    return os.path.normcase(os.path.abspath(a)) == os.path.normcase(os.path.abspath(b))


def _remove_first_slide(prs: "PptxPresentation") -> "SlideLayout":
    """删除演示文稿的第一张幻灯片，返回其 slide_layout 供后续复用。

    说明：python-pptx（经测 1.0.2）未提供公开的删除幻灯片 API，官方社区
    方案即直接操作内部 sldIdLst（见 python-pptx 文档 FAQ “Delete a slide”）。
    此处保留该内部方案，并做三项加固：
      1) 取首张幻灯片布局时，若其 sldId 缺少 r:id 属性，prs.slides[0]
         .slide_layout 会抛 InvalidXmlError，故先捕获并降级；
      2) 即便取得布局，sldId 的 r:id 仍可能为 None，此时 drop_rel(None)
         会抛 KeyError/TypeError，故先判空再删除关系；
      3) 用 try/except 包裹删除过程，若未来版本移除 _sldIdLst 私有属性，
         则优雅降级为“不删除、仅复用布局”，保证生成流程不崩溃
         （最多多一页空白）。
    """
    if len(prs.slides) == 0:
        return prs.slide_layouts[-1]
    # 取布局：rId 缺失会抛 InvalidXmlError，降级为兜底布局且不删除
    try:
        layout: "SlideLayout" = prs.slides[0].slide_layout
    except InvalidXmlError:
        return prs.slide_layouts[-1]
    # _sldIdLst 为 python-pptx 私有属性（官方删除幻灯片方案），用 getattr 兼容未来版本
    sldId_lst = getattr(prs.slides, "_sldIdLst", None)
    if sldId_lst is None:
        return layout
    try:
        sldId = sldId_lst[0]
        rId = sldId.get(qn('r:id'))
        if rId is not None:
            prs.part.drop_rel(rId)
        sldId_lst.remove(sldId)
    except (KeyError, TypeError, InvalidXmlError):
        # 降级：无法安全删除首张幻灯片时，复用已取得的布局（不删除）
        pass
    return layout


def generate_pptx(template_path: str, questions: List[List[str]], font_name: str,
                   font_size: int, output_path: str,
                   line_spacing_type: str = "1 倍", line_spacing_value: float = 1.0,
                   first_line_indent: bool = True,
                   progress_cb: Optional[Callable[[int, int], None]] = None) -> None:
    """基于 PPT 模板为每道题目生成两页幻灯片并保存。

    Raises:
        ValueError: 当 output_path 与 template_path 指向同一文件时，
                    避免覆盖并损坏用户模板。
    """
    # 防御：输出路径与模板路径相同会覆盖并损坏模板文件
    if _same_path(template_path, output_path):
        raise ValueError(
            "输出路径不能与模板路径相同，否则会覆盖并损坏模板文件："
            f"{output_path}"
        )
    prs = Presentation(template_path)
    slide_layout = _remove_first_slide(prs)
    total = len(questions)
    for qi, question in enumerate(questions):
        if progress_cb:
            progress_cb(qi + 1, total)
        for _ in range(2):
            slide = prs.slides.add_slide(slide_layout)
            for shape in list(slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)

            sw = prs.slide_width or 0
            sh = prs.slide_height or 0
            textbox_w = int(int(sw) * 0.8)
            textbox_l = int(int(sw) * 0.1)
            textbox_h = int(int(sh) * 0.75)
            textbox_t = int(int(sh) * 0.15)

            textbox = slide.shapes.add_textbox(textbox_l, textbox_t, textbox_w, textbox_h)
            tf = textbox.text_frame
            tf.word_wrap = True

            for i, para_text in enumerate(question):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                if first_line_indent:
                    para_text = '　　' + para_text

                run = p.add_run()
                run.text = para_text
                _set_font(run, font_name, font_size)
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                p.line_spacing = _resolve_line_spacing(line_spacing_type, line_spacing_value)

    prs.save(output_path)
